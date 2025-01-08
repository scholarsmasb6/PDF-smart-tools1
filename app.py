import os
import re
import nltk
import pandas as pd
import pdfplumber
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer, pipeline
import fitz  # PyMuPDF for extracting images
from flask import Flask, request, jsonify, send_file,render_template
from flask_cors import CORS
import zipfile
import fitz  # PyMuPDF for extracting images
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from nltk.tokenize import sent_tokenize
from pptx.util import Inches, Pt
import PyPDF2


app = Flask(__name__)
CORS(app)
#nltk.data.path.append('C:/Users/allad/AppData/Local/Programs/Python/Python310/nltk_data')
#nltk.download('punkt_tab')

# Initialize directories
UPLOAD_DIR = 'uploads'
OUTPUT_DIR = 'outputs'
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Initialize Models
print("Loading models...")
model_name = "allenai/led-large-16384-arxiv"
model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
tokenizer = AutoTokenizer.from_pretrained(model_name)

summarization_pipeline = pipeline("summarization", model="facebook/bart-large-cnn")

# Utility Functions
def preprocess_text(text):
    """
    Preprocess text to clean and format according to the requirements:
    1. Split sentences based on full stops.
    2. Add newlines for numbered lists or points.
    3. Ensure paragraph separation.
    4. Place headings/subheadings on new lines.
    """
    # Split sentences based on full stops or question marks
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)\s(?=\S)', text)

    # Add newlines for numbered lists or points
    formatted_text = ""
    for sentence in sentences:
        if re.match(r'^\d+\.|\d+\)', sentence.strip()):  # Numbered patterns like 1. or 1)
            formatted_text += f"\n{sentence.strip()}"
        else:
            formatted_text += f" {sentence.strip()}"

    # Split paragraphs based on newlines and ensure paragraph separation
    paragraphs = formatted_text.split("\n")
    formatted_paragraphs = []
    for paragraph in paragraphs:
        if paragraph.strip():
            # Check for headings/subheadings (e.g., uppercase or followed by a colon)
            if paragraph.strip().isupper() or paragraph.strip().endswith(":"):
                formatted_paragraphs.append(f"\n{paragraph.strip()}\n")
            else:
                formatted_paragraphs.append(paragraph.strip())

    return "\n\n".join(formatted_paragraphs)

def extract_text_from_pdf(pdf_path):
    """
    Extract text from a PDF file while meeting the formatting requirements:
    - Avoid mixing table content with main text.
    - Return formatted text.
    """
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                # Extract plain text from the page
                page_text = page.extract_text() or ""

                # Avoid adding table content to plain text
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        table_content = " ".join([" ".join(row) for row in table if row])
                        page_text = page_text.replace(table_content, "")

                text += page_text

        return preprocess_text(text)

    except Exception as e:
        raise Exception(f"Error while extracting text: {e}")



def extract_tables_from_pdf(pdf_path):
    tables = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                extracted_tables = page.extract_tables()
                if extracted_tables:
                    tables.extend(extracted_tables)
        return tables
    except Exception as e:
        raise e

def save_tables_to_excel(tables, output_path):
    try:
        with pd.ExcelWriter(output_path) as writer:
            for i, table in enumerate(tables):
                df = pd.DataFrame(table)
                df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)
        return output_path
    except Exception as e:
        raise e

# Chunk text by tokens
def chunk_text_by_tokens(text, max_tokens=16000):
    sentences = sent_tokenize(text)
    chunks, current_chunk = [], []
    current_chunk_tokens = 0
    for sentence in sentences:
        sentence_tokens = tokenizer.encode(sentence, add_special_tokens=False)
        sentence_token_count = len(sentence_tokens)
        if current_chunk_tokens + sentence_token_count > max_tokens:
            chunks.append(" ".join(current_chunk).strip())
            current_chunk = [sentence]
            current_chunk_tokens = sentence_token_count
        else:
            current_chunk.append(sentence)
            current_chunk_tokens += sentence_token_count
    if current_chunk:
        chunks.append(" ".join(current_chunk).strip())
    return chunks



def extract_important_points(text, num_points=5):
    """
    Extracts important points from summarized text and saves them to a file.
    """
    try:
        if not text.strip():
            raise ValueError("Input text is empty.")

        # Chunk the text into smaller parts for summarization
        chunks = [text[i:i + 16000] for i in range(0, len(text), 16000)]
        summaries = []
        for chunk in chunks:
            # Summarize each chunk, ensuring valid results
            result = summarization_pipeline(chunk, max_length=512, min_length=50, do_sample=False)
            if result and len(result) > 0 and "summary_text" in result[0]:
                summaries.append(result[0]["summary_text"])

        summarized_text = " ".join(summaries)

        # Extract points from the summarized text
        points = summarized_text.split('. ')
        important_points = [f"{i+1}. {point.strip()}." for i, point in enumerate(points) if point.strip()]

        # Limit the points to the specified number
        important_points = important_points[:num_points]

        # Save the points to a file
        output_path = os.path.join(OUTPUT_DIR, "important_points.txt")
        with open(output_path, "w", encoding="utf-8") as file:
            file.write("\n".join(important_points))

        return output_path
    except Exception as e:
        raise Exception(f"Error extracting important points: {e}")


def extract_images_from_pdf(pdf_path, output_dir):
    try:
        os.makedirs(output_dir, exist_ok=True)
        pdf_document = fitz.open(pdf_path)
        image_paths = []

        for i in range(len(pdf_document)):
            page = pdf_document[i]
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_filename = os.path.join(output_dir, f"page_{i+1}_img_{img_index+1}.png")
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes)
                image_paths.append(image_filename)
        
        # Create a ZIP file of the images
        zip_filename = os.path.join(output_dir, "extracted_images.zip")
        with zipfile.ZipFile(zip_filename, 'w') as zipf:
            for img_path in image_paths:
                zipf.write(img_path, os.path.basename(img_path))  # Add image to zip

        return zip_filename  # Return the path to the ZIP file
    except Exception as e:
        raise e



def generate_pdf_presentation(pdf_path, output_dir='outputs'):
    try:
        # Initialize directories
        os.makedirs(output_dir, exist_ok=True)

        # Preprocess and clean the text
        def preprocess_text(text):
            text = re.sub(r'[^\x00-\x7F]+', ' ', text)  # Remove non-ASCII characters
            text = re.sub(r'\s+', ' ', text)  # Remove extra spaces
            return text.strip()

        # Extract text from PDF
        def extract_text_from_pdf(pdf_path):
            text = ""
            try:
                with pdfplumber.open(pdf_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text() or ""
                        text += page_text
                return preprocess_text(text)
            except Exception as e:
                raise e

        # Extract images from PDF
        def extract_images_from_pdf(pdf_path, output_dir):
            try:
                pdf_document = fitz.open(pdf_path)
                image_paths = []
                for i in range(len(pdf_document)):
                    page = pdf_document[i]
                    for img_index, img in enumerate(page.get_images(full=True)):
                        xref = img[0]
                        base_image = pdf_document.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_filename = os.path.join(output_dir, f"page_{i+1}_img_{img_index+1}.png")
                        with open(image_filename, "wb") as img_file:
                            img_file.write(image_bytes)
                        image_paths.append(image_filename)
                return image_paths
            except Exception as e:
                raise e

        # Extract text and images
        text = extract_text_from_pdf(pdf_path)
        images = extract_images_from_pdf(pdf_path, output_dir)

        # Create PowerPoint Presentation
        prs = Presentation()

        # Add title slide
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "PDF to Presentation"
        subtitle.text = "Generated from your PDF"

        # Only add text slide if there is text content
        if text:
            # Add slides for text content
            text_limit = 2000  # Limit of characters per slide
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            
            # Split long text into multiple slides if necessary
            start_index = 0
            while start_index < len(text):
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.shapes.placeholders[1]
                title.text = "Text Content"
                content.text = text[start_index:start_index + text_limit]
                  # Set font size to 12 for the text content
                for paragraph in content.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
                start_index += text_limit

        # Add slides for images only if there are images
        if images:
            for img_path in images:
                slide_layout = prs.slide_layouts[5]  # Blank slide layout
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(8.5), height=Inches(6))

        # Save the presentation
        output_pptx = os.path.join(output_dir, "pdf_presentation.pptx")
        prs.save(output_pptx)
        
        return output_pptx

    except Exception as e:
        raise e

# Function to extract Q&A from PDF
def extract_qa_from_pdf(pdf_path, model_name="deepset/roberta-base-squad2"):
    try:
        # Step 1: Extract text from PDF
        def extract_text(pdf_path):
            text = ""
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return text

        # Step 2: Chunk the text into manageable sizes
        def chunk_text(text, max_tokens=512):
            words = text.split()
            return [" ".join(words[i:i + max_tokens]) for i in range(0, len(words), max_tokens)]

        # Extract the text and chunk it
        text = extract_text(pdf_path)
        text_chunks = chunk_text(text)

        # Step 3: Load the pre-trained QA model
        qa_pipeline = pipeline("question-answering", model=model_name, tokenizer=model_name)

        # Step 4: Define sample questions
        questions = [
            "What is the main topic?",
            "What are the key points?",
            "What details are discussed?",
            "What is the conclusion?",
        ]

        # Step 5: Generate Q&A pairs
        qa_pairs = []
        for chunk in text_chunks:
            for question in questions:
                try:
                    result = qa_pipeline(question=question, context=chunk)
                    qa_pairs.append((question, result.get('answer', 'No answer found')))
                except Exception as e:
                    print(f"Error processing chunk: {e}")

        # Step 6: Save Q&A pairs to a text file
        output_file = "qa_output.txt"
        with open(output_file, "w", encoding="utf-8") as file:
            for i, (question, answer) in enumerate(qa_pairs, 1):
                file.write(f"{i}. Question: {question}\nAnswer: {answer}\n\n")

        return output_file
    except Exception as e:
        print(f"Error extracting Q&A: {e}")
        return None


def merge_pdfs_with_features(pdf1, pdf2, order):
    """
    Merges two PDFs with specific features:
    1. Removes duplicate content.
    2. Maintains headings, subheadings, and structure.
    3. Allows customizable merging order.

    Args:
        pdf1 (file): First PDF file.
        pdf2 (file): Second PDF file.
        order (str): '1' for pdf1 first, '2' for pdf2 first.

    Returns:
        str: Path to the merged PDF file.
    """
    def extract_text(pdf_file):
        """Extract text content from a PDF file."""
        text = ""
        with pdfplumber.open(pdf_file) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text

    # Extract unique content
    content_pdf1 = extract_text(pdf1)
    content_pdf2 = extract_text(pdf2)
    unique_content_pdf2 = "\n".join(
        line for line in content_pdf2.splitlines() if line.strip() and line not in content_pdf1
    )

    pdf_writer = PyPDF2.PdfWriter()

    def add_pages_to_writer(pdf_file, writer):
        reader = PyPDF2.PdfReader(pdf_file)
        for page in reader.pages:
            writer.add_page(page)

    # Merge based on order
    if order == "1":
        add_pages_to_writer(pdf1, pdf_writer)
        add_pages_to_writer(pdf2, pdf_writer)  # Add all pages from PDF2
    elif order == "2":
        add_pages_to_writer(pdf2, pdf_writer)
        add_pages_to_writer(pdf1, pdf_writer)

    output_path = "merged_output.pdf"
    with open(output_path, "wb") as output_file:
        pdf_writer.write(output_file)

    return output_path


# Flask Routes
@app.route('/', methods=["GET"])
def home():
    return render_template('/index.html')


@app.route('/process-pdf', methods=['POST'])
def process_pdf():
    """
    Flask route to process and download the extracted text from an uploaded PDF.
    """
    if 'file' not in request.files:
        return jsonify({"error": "No file provided."}), 400

    pdf_file = request.files['file']
    if pdf_file.filename == '':
        return jsonify({"error": "No selected file."}), 400

    try:
        # Save the uploaded file temporarily
        temp_pdf_path = os.path.join(OUTPUT_DIR, pdf_file.filename)
        pdf_file.save(temp_pdf_path)

        # Extract and process text
        processed_text = extract_text_from_pdf(temp_pdf_path)

        # Save the processed text to a file
        output_file = os.path.join(OUTPUT_DIR, "extracted_text.txt")
        with open(output_file, "w", encoding="utf-8") as file:
            file.write(processed_text)

        # Clean up the temporary PDF file
        os.remove(temp_pdf_path)

        # Send the processed text file for download
        return send_file(output_file, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/process-tables', methods=['POST'])
def process_tables():
    try:
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        # Extract tables from PDF
        tables = extract_tables_from_pdf(pdf_path)
        if not tables:
            return jsonify({"error": "No tables found in the PDF."}), 404

        # Save tables to an Excel file
        output_path = os.path.join(OUTPUT_DIR, f"tables_{file.filename}.xlsx")
        with pd.ExcelWriter(output_path) as writer:
            for i, table in enumerate(tables):
                df = pd.DataFrame(table)
                df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)

        # Send the Excel file as an attachment
        return send_file(output_path, as_attachment=True, download_name=f"tables_{file.filename}.xlsx", mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/summarize', methods=['POST'])
def summarize_pdf():
    try:
        # Get the uploaded file
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file provided"}), 400

        # Save the file in the UPLOAD_DIR
        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        # Extract text from the uploaded PDF
        text = extract_text_from_pdf(pdf_path)
        text = preprocess_text(text)

        # Chunk the text and summarize
        text_chunks = chunk_text_by_tokens(text)
        summaries = []
        for chunk in text_chunks:
            input_ids = tokenizer(chunk, return_tensors="pt", truncation=True).input_ids
            summary_ids = model.generate(input_ids, max_length=512, min_length=50, no_repeat_ngram_size=3)
            summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
            summaries.append(summary)

        summary_text = " ".join(summaries)

        # Save the summary to a file in the UPLOAD_DIR
        summary_file_path = os.path.join(UPLOAD_DIR, "summary.txt")
        with open(summary_file_path, "w", encoding="utf-8") as summary_file:
            summary_file.write(summary_text)

        # Automatically download the file
        return send_file(summary_file_path, as_attachment=True, download_name="summary.txt")

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/extract-images', methods=['POST'])
def extract_images():
    try:
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        output_dir = os.path.join(OUTPUT_DIR, f"images_{file.filename}")
        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        zip_file_path = extract_images_from_pdf(pdf_path, output_dir)
        if not zip_file_path:
            return jsonify({"error": "No images found in the PDF."}), 404

        return send_file(zip_file_path, as_attachment=True, mimetype='application/zip', download_name="extracted_images.zip")

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/extract-points', methods=['POST'])
def extract_points():
    """
    Flask route to upload a PDF, extract important points, and download the result.
    """
    try:
        # Step 1: Validate the uploaded file
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        # Step 2: Save the uploaded PDF
        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        # Step 3: Extract text from the PDF
        text = extract_text_from_pdf(pdf_path)
        if not text:
            return jsonify({"error": "No text found in the PDF."}), 404

        # Step 4: Extract important points from the text
        output_path = extract_important_points(text)

        # Step 5: Return the generated text file for download
        return send_file(output_path, as_attachment=True, mimetype='text/plain', download_name="important_points.txt")

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/upload-pdf', methods=['POST'])
def upload_pdf():
    try:
        # Check if a file is uploaded
        file = request.files.get('file')
        if not file:
            return jsonify({"error": "No file uploaded"}), 400

        # Save the uploaded PDF
        pdf_path = os.path.join(UPLOAD_DIR, file.filename)
        file.save(pdf_path)

        # Generate the PowerPoint presentation
        pptx_path = generate_pdf_presentation(pdf_path)

        # Send the generated PowerPoint presentation for download
        return send_file(pptx_path, as_attachment=True, download_name="pdf_presentation.pptx")

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/extract-q&a', methods=['POST'])
def upload_and_extract_qa():
    try:
        # Check if a file is uploaded
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        if not file.filename.endswith('.pdf'):
            return jsonify({"error": "Only PDF files are allowed"}), 400

        # Save the uploaded file temporarily
        pdf_path = file.filename
        file.save(pdf_path)

        # Extract Q&A and get the output file
        output_file = extract_qa_from_pdf(pdf_path)

        # Send the file as a downloadable response
        if output_file:
            return send_file(output_file, as_attachment=True)
        else:
            return jsonify({"error": "Failed to extract Q&A"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500



@app.route('/merge', methods=['POST'])
def merge_pdfs():
    if 'pdf1' not in request.files or 'pdf2' not in request.files:
        return jsonify({"error": "Both PDF files are required."}), 400

    pdf1 = request.files['pdf1']
    pdf2 = request.files['pdf2']
    order = request.form.get('order', '1')

    try:
        merged_pdf_path = merge_pdfs_with_features(pdf1, pdf2, order)
        return send_file(merged_pdf_path, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    '''finally:
        if os.path.exists(merged_pdf_path):
            os.remove(merged_pdf_path)'''



# Cleanup files periodically (optional)
def cleanup():
    for folder in [UPLOAD_DIR, OUTPUT_DIR]:
        for filename in os.listdir(folder):
            file_path = os.path.join(folder, filename)
            try:
                if os.path.isfile(file_path):
                    os.unlink(file_path)
            except Exception as e:
                print(f"Failed to delete {file_path}. Reason: {e}")

# Main Driver
if __name__ == '__main__':
    app.run(debug=False, port=5000)




